import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import anthropic
import base64
import json
import fitz
from dotenv import load_dotenv
import os
import re

# ------ GRAPH MAKER IMPORTS ------
import matplotlib
matplotlib.use("Agg")  # Needed for Streamlit / headless environments
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np

# ---------------- GRAPH MAKER FUNCTIONS ----------------

def create_percentage_donut(percentage, title, colors=('#2596be', '#f0f0f0')):
    fig, ax = plt.subplots()
    ax.pie(
        [percentage, 100 - percentage],
        colors=colors,
        startangle=90,
        counterclock=False,
        wedgeprops={'width': 0.3, 'edgecolor': 'white'}
    )
    ax.text(
        0, 0, f'{percentage}%',
        horizontalalignment='center',
        verticalalignment='center',
        fontsize=24,
        fontweight='bold'
    )
    ax.text(
        0, -0.1, title,
        horizontalalignment='center',
        fontsize=12,
        wrap=True,
        transform=ax.transAxes
    )
    ax.set(aspect="equal")
    plt.tight_layout()
    return fig, ax

def create_patient_count_visualization(count, title, icons_per_row=35):
    fig, ax = plt.subplots()
    rows = int(np.ceil(count / icons_per_row))
    x_coords = []
    y_coords = []
    for i in range(count):
        row = i // icons_per_row
        col = i % icons_per_row
        x_coords.append(col)
        y_coords.append(-row)
    ax.plot(
        x_coords,
        y_coords,
        linestyle='none',
        marker='$\\bigodot$',
        markersize=12,
        color='#2596be',
        markeredgewidth=0
    )
    ax.text(
        icons_per_row / 2,
        1,
        str(count),
        horizontalalignment='center',
        verticalalignment='bottom',
        fontsize=36,
        fontweight='bold'
    )
    ax.text(
        icons_per_row / 2,
        0.6,
        title,
        horizontalalignment='center',
        verticalalignment='top',
        fontsize=14,
        color='#2596be'
    )
    ax.set_xlim(-1, icons_per_row + 1)
    ax.set_ylim(-rows - 1, 2)
    ax.axis('off')
    plt.tight_layout()
    return fig, ax

def create_comparison_bars(values, labels, title, color='#2596be'):
    fig, ax = plt.subplots()
    y_pos = np.arange(len(values))
    ax.barh(y_pos, values, height=0.6, color=color)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels)
    ax.set_xlabel('Value')
    ax.set_title(title, pad=15, fontsize=14)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    return fig, ax

def create_trend_line(values, dates, title, color='#2596be'):
    fig, ax = plt.subplots()
    ax.plot(
        dates,
        values,
        color=color,
        linewidth=2,
        marker='o',
        markersize=6,
        markerfacecolor='white',
        markeredgewidth=2
    )
    ax.grid(True, linestyle='--', alpha=0.7)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title(title, pad=15, fontsize=14)
    plt.xticks(rotation=45)
    plt.tight_layout()
    return fig, ax

def create_stacked_percentage(categories, percentages, title, colors=None):
    fig, ax = plt.subplots()
    if colors is None:
        colors = ['#2596be', '#1a7fa3', '#0f6788', '#044e6d', '#003652']
        colors = colors[:len(categories)]
    left = 0
    for i, percentage in enumerate(percentages):
        ax.barh(
            0,
            percentage,
            left=left,
            height=0.5,
            color=colors[i],
            label=f'{categories[i]} ({percentage}%)'
        )
        if percentage > 5:
            ax.text(
                left + percentage / 2,
                0,
                f'{percentage}%',
                horizontalalignment='center',
                verticalalignment='center',
                color='white',
                fontweight='bold'
            )
        left += percentage
    ax.set_xlim(0, 100)
    ax.set_ylim(-0.6, 0.6)
    ax.set_title(title, pad=15, fontsize=14)
    ax.axis('off')
    ax.legend(
        bbox_to_anchor=(0.5, -0.5),
        loc='upper center',
        ncol=len(categories),
        frameon=False
    )
    plt.tight_layout()
    return fig, ax

def fig_to_image_stream(fig):
    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream

# ---------------- END GRAPH MAKER FUNCTIONS ----------------

# Load environment variables
load_dotenv()

st.set_page_config(page_title="PDF to PowerPoint Generator", layout="wide")

api_key = os.getenv("ANTHROPIC_API_KEY")
if not api_key:
    st.error("Anthropic API key not found. Please check your .env file.")
    st.stop()

client = anthropic.Anthropic(
    api_key=api_key,
    default_headers={
        "anthropic-beta": "pdfs-2024-09-25"
    }
)

# Model Name
MODEL_NAME = "claude-3-5-sonnet-20241022"  # <-- as requested

GRAPH_PROMPT_PATH = "prompt2.txt"  # <-- your prompt file

@st.cache_resource
def load_prompt_text(prompt_path):
    with open(prompt_path, "r") as file:
        return file.read()

def display_logo():
    logo_svg = '''
    <div style="text-align: center; margin-top: 2rem;">
        <svg viewBox="0 0 200 80" style="width: 150px;">
            <text x="10" y="60" font-family="Arial" font-size="60" font-weight="bold" fill="#FFFFFF">
                Peel
            </text>
        </svg>
    </div>
    '''
    st.markdown(logo_svg, unsafe_allow_html=True)

def extract_text_from_pdf(pdf_bytes):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def generate_json_using_claude(prompt, pdf_bytes, simplification_level, progress_callback=None):
    """Sends the prompt and PDF data to Claude to generate JSON data."""
    formatted_prompt = prompt.replace("{{SIMPLIFICATION_LEVEL}}", str(simplification_level))
    pdf_data = base64.b64encode(pdf_bytes).decode('utf-8')

    # Prefix the prompt with "\n\nHuman:" to comply with Claude's requirements
    prefixed_prompt = "\n\nHuman:" + formatted_prompt

    messages = [
        {
            "role": 'user',
            "content": [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": pdf_data
                    }
                },
                {
                    "type": "text",
                    "text": prefixed_prompt
                }
            ]
        }
    ]

    try:
        if progress_callback:
            progress_callback("Information peeling", 10)
        response = client.messages.create(
            model=MODEL_NAME,
            max_tokens=8192,
            messages=messages,
        )

        if progress_callback:
            progress_callback("Information peeling.", 70)

        assistant_reply = response.content[0].text
        json_output_pattern = r"<json_output>\s*(\{.*?\})\s*</json_output>"
        match = re.search(json_output_pattern, assistant_reply, re.DOTALL)

        if match:
            json_string = match.group(1)
            json_data = json.loads(json_string)
            
            # Ensure exactly three charts
            charts = json_data.get("CHARTS", [])
            if len(charts) < 3:
                # Fill with None or default charts if fewer than 3
                for _ in range(3 - len(charts)):
                    charts.append(None)
            elif len(charts) > 3:
                charts = charts[:3]
            json_data["CHARTS"] = charts

            return json_data
        else:
            try:
                json_data = json.loads(assistant_reply)
                # Ensure exactly three charts
                charts = json_data.get("CHARTS", [])
                if len(charts) < 3:
                    for _ in range(3 - len(charts)):
                        charts.append(None)
                elif len(charts) > 3:
                    charts = charts[:3]
                json_data["CHARTS"] = charts
                return json_data
            except json.JSONDecodeError:
                st.error("Failed to find or parse JSON output in the response.")
                return {}
    except Exception as e:
        st.error(f"Error communicating with Claude: {e}")
        return {}
    finally:
        if progress_callback:
            progress_callback("Completed.", 100)

def populate_ppt_template(json_data, prs):
    """
    Insert text into existing slides placeholders AND replace graph placeholders with charts.
    This function dynamically detects graph placeholders, retrieves their dimensions,
    generates charts accordingly, and inserts them into the placeholders.
    """
    # 1) Insert text into placeholders
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("(Title)", json_data.get("Title", ""))
                        run.text = run.text.replace("(AUTHOR_NAMES)", ", ".join(json_data.get("AUTHOR_NAMES", [])))
                        run.text = run.text.replace("(PAPER_PMID)", json_data.get("PAPER_PMID", ""))
                        run.text = run.text.replace("(PAPER_DOI)", json_data.get("PAPER_DOI", ""))
                        run.text = run.text.replace("(Background_Info)", json_data.get("Background_Info", ""))
                        run.text = run.text.replace("(Patient Quote)", json_data.get("Patient_Quote", ""))
                        run.text = run.text.replace("(Patient name)", json_data.get("Patient_Name", ""))
                        run.text = run.text.replace("(AIMS)", "\n".join(json_data.get("AIMS", [])))
                        run.text = run.text.replace("(Methods)", json_data.get("Methods", ""))
                        run.text = run.text.replace("(Findings)", "\n".join(json_data.get("Findings", [])))
                        run.text = run.text.replace("(Conclusion)", json_data.get("Conclusion", ""))
                        run.text = run.text.replace("(Slide_Number)", str(slide_number))

    # 2) Replace graph placeholders with charts
    charts_list = json_data.get("CHARTS", [])
    if not charts_list:
        st.warning("No charts found in JSON data.")
        return

    # Extract all graph placeholders across all slides
    graph_placeholders = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name.startswith("graph_"):
                graph_placeholders.append(shape)

    num_placeholders = len(graph_placeholders)
    num_charts = len(charts_list)

    if num_charts < num_placeholders:
        st.warning(f"Number of charts ({num_charts}) is less than the number of placeholders ({num_placeholders}). Some placeholders will remain empty.")
    elif num_charts > num_placeholders:
        st.warning(f"Number of charts ({num_charts}) is greater than the number of placeholders ({num_placeholders}). Extra charts will be ignored.")

    # Determine the number of charts to process
    num_to_process = min(num_charts, num_placeholders)

    for idx in range(num_to_process):
        chart_info = charts_list[idx]
        placeholder_shape = graph_placeholders[idx]
        placeholder_name = placeholder_shape.name

        if not chart_info:
            st.warning(f"Chart #{idx+1} data is missing. Placeholder '{placeholder_name}' will not be replaced.")
            continue  # Skip if no chart info

        chart_type = chart_info.get("chart_type", "")
        chart_title = chart_info.get("chart_title", "")
        data = chart_info.get("data", {})

        fig, ax = None, None
        if chart_type == "donut":
            percentage = data.get("percentage", 50)
            fig, ax = create_percentage_donut(percentage, chart_title)
        elif chart_type == "patient_count":
            count = data.get("count", 0)
            fig, ax = create_patient_count_visualization(count, chart_title)
        elif chart_type == "comparison_bars":
            values = data.get("values", [])
            labels = data.get("labels", [])
            fig, ax = create_comparison_bars(values, labels, chart_title)
        elif chart_type == "trend_line":
            values = data.get("values", [])
            dates = data.get("dates", [])
            fig, ax = create_trend_line(values, dates, chart_title)
        elif chart_type == "stacked_percentage":
            categories = data.get("categories", [])
            percentages = data.get("percentages", [])
            fig, ax = create_stacked_percentage(categories, percentages, chart_title)
        else:
            st.warning(f"Chart type '{chart_type}' is not recognized. Skipping chart #{idx+1}.")
            continue

        if fig:
            # Retrieve placeholder dimensions in inches
            width_in = placeholder_shape.width.inches
            height_in = placeholder_shape.height.inches

            # Adjust the figure size to match placeholder dimensions
            fig.set_size_inches(width_in, height_in)

            # Save the figure to an image stream
            pic_stream = fig_to_image_stream(fig)

            # Insert the image into the slide
            slide = placeholder_shape.part.slide
            left = placeholder_shape.left
            top = placeholder_shape.top
            width = placeholder_shape.width
            height = placeholder_shape.height

            # Add the picture
            pic = slide.shapes.add_picture(pic_stream, left, top, width=width, height=height)

            # Remove the placeholder shape
            sp = placeholder_shape._element
            sp.getparent().remove(sp)

            st.info(f"Replaced placeholder '{placeholder_name}' with chart #{idx+1}: {chart_type}")

            # Close the figure to free memory
            plt.close(fig)

    # Handle any extra charts (if charts > placeholders)
    if num_charts > num_placeholders:
        for idx in range(num_placeholders, num_charts):
            st.warning(f"Chart #{idx+1} has no corresponding placeholder and will be ignored.")

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def edit_json_section(json_data, section_name):
    # Let user edit each section in Streamlit
    if section_name == "Paper Information":
        st.text_input("Title", value=json_data.get("Title", ""), key="title")
        authors_str = ", ".join(json_data.get("AUTHOR_NAMES", []))
        authors_str = st.text_input("Authors (comma-separated)", value=authors_str, key="authors")
        st.text_input("PMID", value=json_data.get("PAPER_PMID", ""), key="pmid")
        st.text_input("DOI", value=json_data.get("PAPER_DOI", ""), key="doi")

        json_data["Title"] = st.session_state.title
        json_data["AUTHOR_NAMES"] = [a.strip() for a in st.session_state.authors.split(",") if a.strip()]
        json_data["PAPER_PMID"] = st.session_state.pmid
        json_data["PAPER_DOI"] = st.session_state.doi

    elif section_name == "Background":
        st.text_area("Background Info", value=json_data.get("Background_Info", ""), key="background")
        json_data["Background_Info"] = st.session_state.background

    elif section_name == "Patient Information":
        st.text_area("Patient Quote", value=json_data.get("Patient_Quote", ""), key="patient_quote")
        st.text_input("Patient Name", value=json_data.get("Patient_Name", ""), key="patient_name")
        json_data["Patient_Quote"] = st.session_state.patient_quote
        json_data["Patient_Name"] = st.session_state.patient_name

    elif section_name == "Research Details":
        st.text_area("Aims (one per line)", value="\n".join(json_data.get("AIMS", [])), key="aims")
        st.text_area("Methods", value=json_data.get("Methods", ""), key="methods")
        st.text_area("Findings (one per line)", value="\n".join(json_data.get("Findings", [])), key="findings")
        st.text_area("Conclusion", value=json_data.get("Conclusion", ""), key="conclusion")

        json_data["AIMS"] = [a.strip() for a in st.session_state.aims.split("\n") if a.strip()]
        json_data["Methods"] = st.session_state.methods
        json_data["Findings"] = [f.strip() for f in st.session_state.findings.split("\n") if f.strip()]
        json_data["Conclusion"] = st.session_state.conclusion

    return json_data

@st.cache_resource
def load_ppt_template(template_path):
    return Presentation(template_path)

def show_chart_previews(json_data):
    """
    Create & show all charts in Streamlit, if any, using st.pyplot(fig).
    """
    charts = json_data.get("CHARTS", [])
    if not charts:
        st.info("No charts found in JSON.")
        return

    for i, chart_info in enumerate(charts, start=1):
        if not chart_info:
            st.warning(f"Chart #{i} is missing.")
            continue

        chart_type = chart_info.get("chart_type", "")
        chart_title = chart_info.get("chart_title", "")
        data = chart_info.get("data", {})

        fig = None
        if chart_type == "donut":
            fig, _ = create_percentage_donut(data.get("percentage", 50), chart_title)
        elif chart_type == "patient_count":
            fig, _ = create_patient_count_visualization(data.get("count", 0), chart_title)
        elif chart_type == "comparison_bars":
            fig, _ = create_comparison_bars(data.get("values", []), data.get("labels", []), chart_title)
        elif chart_type == "trend_line":
            fig, _ = create_trend_line(data.get("values", []), data.get("dates", []), chart_title)
        elif chart_type == "stacked_percentage":
            fig, _ = create_stacked_percentage(data.get("categories", []), data.get("percentages", []), chart_title)
        else:
            st.warning(f"Chart type '{chart_type}' is not recognized. Skipping chart #{i}.")
            continue

        if fig is not None:
            st.markdown(f"#### Chart Preview #{i}: {chart_type}")
            st.pyplot(fig)
        else:
            st.warning(f"Chart #{i} could not be created.")

def main():
    display_logo()
    st.markdown('<h1 class="main-title">PDF to PowerPoint Generator</h1>', unsafe_allow_html=True)
    st.markdown("""
        <div class="intro-text">
            Peel back the complexity, let your research speak to everyone.
        </div>
    """, unsafe_allow_html=True)

    # Session states
    if "json_data" not in st.session_state:
        st.session_state.json_data = None
    if "ppt_file" not in st.session_state:
        st.session_state.ppt_file = None
    if "current_section" not in st.session_state:
        st.session_state.current_section = "Paper Information"

    # Load your prompt from graphprompt.txt
    prompt_text = load_prompt_text(GRAPH_PROMPT_PATH)
    TEMPLATE_PATH = "Template.pptx"

    uploaded_pdf = st.file_uploader("Upload your PDF file", type="pdf")

    # Simplification level
    simplification_level = st.select_slider(
        "Select Simplification Level",
        options=list(range(1, 11)),
        format_func=lambda x: "Academic" if x == 1 else "Patient" if x == 10 else f"Level {x}",
        value=5
    )

    if uploaded_pdf:
        pdf_bytes = uploaded_pdf.read()

        if st.button("Let’s Peel"):
            with st.spinner("Information peeling"):
                progress_bar = st.progress(0)
                status_text = st.empty()

                def update_progress(step_description, pct):
                    progress_bar.progress(pct)
                    status_text.text(step_description)

                # Call Claude with the PDF + prompt
                st.session_state.json_data = generate_json_using_claude(
                    prompt_text,
                    pdf_bytes,
                    simplification_level,
                    progress_callback=update_progress
                )
                if st.session_state.json_data:
                    st.success("JSON data generated successfully!")

    # If we have JSON in the session
    if st.session_state.json_data:
        sections = ["Paper Information", "Background", "Patient Information", "Research Details"]
        st.session_state.current_section = st.radio(
            "Select section to edit:",
            sections,
            key="section_selector"
        )

        # Let user edit
        st.session_state.json_data = edit_json_section(
            st.session_state.json_data,
            st.session_state.current_section
        )

        # Optionally show entire JSON (including CHARTS)
        if st.checkbox("Show Current JSON Data"):
            st.json(st.session_state.json_data)

        # Let user preview charts in Streamlit
        if st.checkbox("Show Chart Previews"):
            show_chart_previews(st.session_state.json_data)

        # Finally, generate PPT
        if st.button("Generate PowerPoint"):
            with st.spinner("Generating PowerPoint..."):
                prs = load_ppt_template(TEMPLATE_PATH)
                st.session_state.ppt_file = populate_ppt_template(st.session_state.json_data, prs)
            st.success("PowerPoint generated successfully!")

    # If PPT is generated, show download button
    if st.session_state.ppt_file:
        st.download_button(
            label="Download PowerPoint",
            data=st.session_state.ppt_file,
            file_name="presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

if __name__ == "__main__":
    main()
